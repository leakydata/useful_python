########## LIST OF ALL IMPORTS FOR EVERY FUNCTION ##########
import re
import numpy as np 
import pandas as pd

from pptx import Presentation
from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE, MSO_THEME_COLOR ##MSO_THEME_COLOR https://python-pptx.readthedocs.io/en/latest/api/enum/MsoThemeColorIndex.html
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml

from collections import Counter

from bs4 import BeautifulSoup

## -----------------------------------------------------------

#Check Shape Type in a slide or slides
for shape in slide.shapes:
  print(shape.shape_type)
  
  
########## CHARTS ##########


## Axes Values
# Maximum Y-Axis Value (the maximun scale)
prs.slides[40].shapes[0].chart.plots[0].chart.value_axis.maximum_scale

# Minimum Y-Axis Value (the minimum scale)
prs.slides[40].shapes[0].chart.plots[0].chart.value_axis.minimum_scale

# X-Axis Labels
tupleOfTuples = prs.slides[40].shapes[0].chart.plots[0].categories.flattened_labels
x_axis_labels = list(sum(tupleOfTuples, ()))
print(x_axis_labels)

# Does the plot axis have Minor Gridlines (boolean: True or False)
prs.slides[40].shapes[0].chart.plots[0].chart.value_axis.has_minor_gridlines

# Does the plot axis have Major Gridlines (boolean: True or False)
prs.slides[40].shapes[0].chart.plots[0].chart.value_axis.has_major_gridlines


## Plot Values
# The chart plot bargraph values
prs.slides[40].shapes[0].chart.series[0].values

# The chart plot values for a stacked bargraph
for s in range(0,len(prs.slides[40].shapes[1].chart.series)):
    print(prs.slides[40].shapes[1].chart.series[s].values)

# Check if chart has legend and list the values
legend = []
if prs.slides[40].shapes[1].chart.has_legend:
    for p in range(0, len(prs.slides[40].shapes[1].chart.plots)):
        for s in range(0, len(prs.slides[40].shapes[1].chart.plots[p].series)):
            print(prs.slides[40].shapes[1].chart.plots[p].series[s].name)
            legend.append(prs.slides[40].shapes[1].chart.plots[p].series[s].name)
    
    
#----------------
# Chart Type
prs.slides[40].shapes[1].chart.chart_type # COLUMN_STACKED (52)
prs.slides[40].shapes[1].chart.chart_type._member_name # COLUMN_STACKED
prs.slides[40].shapes[1].chart.chart_type.real # 52
prs.slides[40].shapes[1].chart.chart_type._docstring # Stacked Column.

# Pull all the values out of charts and convert them into dicts then dataframes
from collections import defaultdict
for slide_idx,slide in enumerate(prs.slides):
    for i,shape in enumerate(slide.shapes):
        chart_data = defaultdict(list) 
        
        if shape.shape_type == MSO_SHAPE_TYPE.CHART: # THIS IS WHERE THE CHART SHAPE STARTS (All other shapes are excluded for now)
            
            ch = shape.chart 
            ch_labels = x_labels(ch)
            ch_series = chart_vals(ch)
            chart_name = shape.name.upper()
            
            print("\n\n----------------- SLIDE:",slide_idx+1,"SHAPE:",i, chart_name,"------------------------------")
            
            if len(ch_series) > 0:
                empty_column_indexes = []
                for idx,l in enumerate(ch_series):
                    if all(v is None for v in l):
                        empty_column_indexes.append(idx)

                for s in range(len(ch_series)):
                    for key, value in zip(ch_labels,ch_series[s]):
                        chart_data[key].append(value)
                
                if '' in chart_data:
                    del chart_data['']
                    
                df = pd.DataFrame.from_dict(chart_data)
                df.dropna(axis=0,how='all',inplace=True) 
                final_chart_data = list(df.T.to_dict().values())
                
                legend = []
                if ch.has_legend:
                    legend = get_legend(ch)
                    if bool(empty_column_indexes):
                        for del_col_idx in empty_column_indexes:
                            del legend[del_col_idx]
                    final_chart_data = dict(zip(legend,final_chart_data))
                    
                else:        
                    if bool(final_chart_data):
                        if bool(ch.plots[0].series[0]._ser.xpath('//c:tx//c:v/text()')):
                            legend = list(ch.plots[0].series[0]._ser.xpath('//c:tx//c:v/text()'))
                            if bool(empty_column_indexes):
                                for del_col_idx in empty_column_indexes:
                                    del legend[del_col_idx]
                            final_chart_data = dict(zip(legend, final_chart_data))

                        if type(final_chart_data) is list:
                            final_chart_data = final_chart_data[0]
                   
                    if dict_levels(final_chart_data) > 1:
                        print(pd.DataFrame.from_dict(final_chart_data))
                    else:
                        ind = [str(x) for x in range(len(final_chart_data))]
                        final_chart_data = dict(zip(ind, [final_chart_data]))
                        print(pd.DataFrame.from_dict(final_chart_data))
            else:
                print("SLIDE#:",slide_idx+1, "FINAL DATA:", "-- EMPTY CHART --")


#------- Smart Shapes i.e. Diagrams ----------

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml

# Convert theme colors into hex color codes
def theme_color_convert(slide_num,color):
    slide = prs.slides[slide_num]

    slide_master = slide.slide_layout.slide_master
    slide_master_part = slide_master.part
    theme_part = slide_master_part.part_related_by(RT.THEME)
    theme = parse_xml(theme_part.blob)  # theme here is an <a:theme> element
    
    if 'dk1' in color or 'lt1' in color:
        x = "a:themeElements/a:clrScheme/a:"+color+"/a:sysClr/@lastClr"
    elif 'bg' in color or 'tx' in color:
        x = "//@"+color
    else:
        x = "a:themeElements/a:clrScheme/a:"+color+"/a:srgbClr/@val"
    
    rgb_color = theme.xpath(x)[0]
    
    while('bg' in rgb_color or 'tx' in rgb_color or 'lt' in rgb_color or 'dk' in rgb_color):
        rgb_color = theme_color_convert(slide_num,rgb_color)
    return rgb_color

from bs4 import BeautifulSoup
# Get the rIds of chevrons in Diagram (basically where text is likely to be stored
def get_rIds(slide):
    soup = BeautifulSoup(prs.slides[slide].shapes[0].part.rels.xml)
    #print(soup)
    diagram_rIds = []
    text_data = soup.findAll('relationship', attrs={'target' : re.compile(r'(^\.\./diagrams/drawing[0-9]*\.xml$)')}) 
    #print(text_data)
    for t in text_data:
        diagram_rIds.append(t.get('id'))
    return diagram_rIds
    
# Once you have the text about pull out text and combine subfields (specific to my own work and digrams: we use numbers to label)       
def diagram_text(slide,shape_num,rId):
    data = prs.slides[slide].shapes[shape_num].part.related_parts[rId].blob
    soup = BeautifulSoup(data)
    text_data = soup.findAll('a:t') 

    mydict = {}
    new_list = [n.contents[0] if type(n.contents[0]) != 'int' else int(n.contents[0]) for n in text_data]
    digits = [int(idx+1) for idx, val in enumerate(new_list) if val.isdigit()]

    for idx, i in enumerate(digits):
        if idx != len(digits)-1:
            mydict[new_list[i-1]] = ''.join(new_list[i:i+1])
        else:
            mydict[new_list[i-1]] = ''.join(new_list[i:])
    return mydict

from collections import Counter
def get_highlight_color(slide,shape_num,rId):
    data = prs.slides[slide].shapes[shape_num].part.related_parts[rId].blob
    soup = BeautifulSoup(data)
    shapes_data = soup.findAll('dsp:sp') 

    colors = []
    for sh in shapes_data:
        for t in sh.find('a:solidfill'):
            colors.append(t.get('val'))
    print(colors)
    print(colors[1::2]) #start with index 1 #FFFFFF and return every other element
    box_colors = colors[1::2]
    if 'FFFFFF' in box_colors:
        box_colors.remove('FFFFFF')
    return min(Counter(box_colors), key=Counter(box_colors).get)

# Get the highlighted chevron text box in the diagram/smart art
def highlighted_box(slide,shape_num,rId):
    data = prs.slides[slide].shapes[shape_num].part.related_parts[rId].blob
    soup = BeautifulSoup(data)
    shapes_data = soup.findAll('dsp:sp') 

    colors = []
    
    for sh in shapes_data:
        for t in sh.find('a:solidfill'):
            colors.append(t.get('val'))
            
    fixed_colors = []
    for c in colors:
        pattern = re.compile("[A-Za-z][A-Za-z][A-Za-z]*[0-9]*")
        if pattern.match(c):
            f_color = theme_color_convert(slide,c)
            fixed_colors.append(f_color)
        else:
            fixed_colors.append(c)

    print(fixed_colors[1::2]) #start with index 1 in this one case #FFFFFF and return every other element

    box_colors = list(filter(lambda a: a != 'FFFFFF', fixed_colors[1::2])) #remove all #ffffff from the list
    print("BOX COLORS:",box_colors)

    if len(Counter(box_colors)) > 0:
        highlight_color = min(Counter(box_colors), key=Counter(box_colors).get)
        print("HIGHLIGHT COLOR:",highlight_color)
        return fixed_colors[1::2].index(highlight_color)
    else:
        return -1
    
